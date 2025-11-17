"""
PPTX to text converter.
"""
from pptx import Presentation
from .base import BaseConverter


class PPTXConverter(BaseConverter):
    """Converter for PPTX presentations."""

    def convert(self, input_path: str, output_path: str) -> bool:
        """
        Convert PPTX to text.

        Args:
            input_path: Path to PPTX file
            output_path: Path to output text file

        Returns:
            True if successful, False otherwise
        """
        try:
            text = self._extract_text(input_path)
            if text:
                return self._write_output(text, output_path)
            return False
        except Exception as e:
            self.logger.error(f"Error converting PPTX {input_path}: {str(e)}")
            return False

    def _extract_text(self, input_path: str) -> str:
        """
        Extract text from PPTX.

        Args:
            input_path: Path to PPTX file

        Returns:
            Extracted text
        """
        prs = Presentation(input_path)
        text_parts = []

        for slide_num, slide in enumerate(prs.slides, 1):
            if self.preserve_structure:
                text_parts.append(f"\n{'='*60}\n")
                text_parts.append(f"Slide {slide_num}\n")
                text_parts.append(f"{'='*60}\n\n")

            # Extract text from shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text = shape.text.strip()
                    if text:
                        text_parts.append(f"{text}\n")

                # Extract text from tables
                if shape.has_table:
                    if self.preserve_structure:
                        text_parts.append("\n[TABLE]\n")

                    for row in shape.table.rows:
                        row_text = []
                        for cell in row.cells:
                            cell_text = cell.text.strip()
                            if cell_text:
                                row_text.append(cell_text)

                        if row_text:
                            if self.preserve_structure:
                                text_parts.append(" | ".join(row_text) + "\n")
                            else:
                                text_parts.append(" ".join(row_text) + "\n")

                    if self.preserve_structure:
                        text_parts.append("[/TABLE]\n")

            text_parts.append("\n")

        return ''.join(text_parts).strip()

    @classmethod
    def get_supported_extensions(cls) -> list[str]:
        """Get supported file extensions."""
        return ['.pptx']
